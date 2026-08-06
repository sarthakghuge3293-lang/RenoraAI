import os
from services.document_reader import DocumentReader

def create_test_files():
    # Excel
    import pandas as pd
    df = pd.DataFrame({'Name': ['Alice', 'Bob'], 'Salary': [50000, 60000]})
    df.to_excel('test.xlsx', index=False)
    
    # CSV
    df.to_csv('test.csv', index=False)
    
    # Word
    from docx import Document
    doc = Document()
    doc.add_heading('Test Word', 0)
    doc.add_paragraph('This is a test paragraph.')
    doc.save('test.docx')
    
    # PPTX
    from pptx import Presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Hello, World!"
    prs.save('test.pptx')

def test_reader():
    reader = DocumentReader()
    
    files = ['test.xlsx', 'test.csv', 'test.docx', 'test.pptx']
    for f in files:
        print(f"Testing {f}...")
        try:
            pages = reader.read_document(f, f)
            print(f"Success! Pages: {len(pages)}, First page preview: {pages[0]['text'][:50]}")
        except Exception as e:
            print(f"Failed {f}: {str(e)}")

if __name__ == '__main__':
    create_test_files()
    test_reader()
