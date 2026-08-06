import os
import pandas as pd
from docx import Document
from pptx import Presentation
from typing import List, Dict

# Import existing PDF reader
from services.pdf_reader import PDFReader


class DocumentReader:
    def __init__(self):
        self.pdf_reader = PDFReader()

    def read_document(self, file_path: str, original_name: str = None) -> List[Dict]:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")

        # Determine file extension
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self.pdf_reader.read_pdf(file_path)
        elif ext in [".xlsx", ".xls"]:
            return self.read_excel(file_path)
        elif ext == ".csv":
            return self.read_csv(file_path)
        elif ext == ".docx":
            return self.read_docx(file_path)
        elif ext == ".pptx":
            return self.read_pptx(file_path)
        else:
            raise ValueError(f"Unsupported document type: {ext}")

    def read_excel(self, file_path: str) -> List[Dict]:
        """
        Reads all sheets in an Excel file.
        Each sheet is treated as one 'page'.
        """
        pages = []
        try:
            excel_data = pd.read_excel(file_path, sheet_name=None)
            sheet_index = 1
            for sheet_name, df in excel_data.items():
                # Drop rows/columns that are completely empty
                df = df.dropna(how="all").dropna(axis=1, how="all")

                # Convert DataFrame to a string representation
                text = f"Sheet Name: {sheet_name}\n"

                if df.empty:
                    text += "Empty Sheet"
                else:
                    text += df.to_string(index=False)

                pages.append(
                    {
                        "page": sheet_index,
                        "text": text.strip(),
                        "word_count": len(text.split()),
                        "char_count": len(text),
                    }
                )
                sheet_index += 1
        except Exception as e:
            raise RuntimeError(f"Failed to read Excel file: {str(e)}")

        return pages

    def read_csv(self, file_path: str) -> List[Dict]:
        """
        Reads a CSV file. Treats the entire CSV as a single 'page' or chunk it into multiple
        if it's very large, but for now treating as one page.
        """
        pages = []
        try:
            df = pd.read_csv(file_path)
            # Drop empty
            df = df.dropna(how="all").dropna(axis=1, how="all")

            text = df.to_string(index=False)

            pages.append(
                {
                    "page": 1,
                    "text": text.strip(),
                    "word_count": len(text.split()),
                    "char_count": len(text),
                }
            )
        except Exception as e:
            raise RuntimeError(f"Failed to read CSV file: {str(e)}")

        return pages

    def read_docx(self, file_path: str) -> List[Dict]:
        """
        Reads a Word document. Since Word doesn't have strict 'pages',
        we treat the whole document as page 1, or we could split by sections.
        Treating as page 1 for simplicity and let the chunker handle the rest.
        """
        pages = []
        try:
            doc = Document(file_path)
            text_blocks = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_blocks.append(para.text.strip())

            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    text_blocks.append(" | ".join(row_data))

            text = "\n".join(text_blocks)

            pages.append(
                {
                    "page": 1,
                    "text": text.strip(),
                    "word_count": len(text.split()),
                    "char_count": len(text),
                }
            )
        except Exception as e:
            raise RuntimeError(f"Failed to read Word document: {str(e)}")

        return pages

    def read_pptx(self, file_path: str) -> List[Dict]:
        """
        Reads a PowerPoint file.
        Each slide is treated as a 'page'.
        """
        pages = []
        try:
            prs = Presentation(file_path)
            slide_index = 1

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # Extract notes if any
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_text.append(f"Notes: {notes}")

                text = "\n".join(slide_text)

                pages.append(
                    {
                        "page": slide_index,
                        "text": text.strip(),
                        "word_count": len(text.split()),
                        "char_count": len(text),
                    }
                )
                slide_index += 1

        except Exception as e:
            raise RuntimeError(f"Failed to read PowerPoint file: {str(e)}")

        return pages
